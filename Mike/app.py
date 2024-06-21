import streamlit as st
import os

from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings 
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.chains.question_answering import load_qa_chain

####### Funciones para leer cualquier tipo de documento #######

# DOCX
import docx #pip install python-docx 
def docx_list(nombre):
    document = docx.Document(nombre)
    texto = ''
    for parrafo in document.paragraphs:
        texto+=parrafo.text + ' '
    lista = [texto]
    lista = [elemento.replace('\n', '').replace('\t', '') for elemento in lista]
    return lista

# PDF
import PyPDF2 #pip install PyPDF2
def pdf_list(nombre):
    texto = ''
    with open(nombre, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_paginas = len(pdf_reader.pages)
        for pagina in range(num_paginas):
            texto += pdf_reader.pages[pagina].extract_text() + ' '
    lista = [texto]
    lista = [elemento.replace('\n', '').replace('\t', '') for elemento in lista]
    return lista

# PPTX
from pptx import Presentation #pip install python-pptx
def pptx_list(nombre):
    texto = ''
    presentation = Presentation(nombre)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texto += shape.text + ' '
    lista = [texto]
    lista = [elemento.replace('\n', '').replace('\t', '') for elemento in lista]
    return lista

# XLSX
import openpyxl #pip install openpyxl
def xlsx_list(nombre):
    texto = ''
    wb = openpyxl.load_workbook(nombre)
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    texto += str(cell) + ' '
    lista = [texto]
    lista = [elemento.replace('\n', '').replace('\t', '') for elemento in lista]
    return lista

# TXT
def txt_list(nombre):
    with open(nombre, 'r', encoding='utf-8') as file:
        texto = file.read()
    lista = [texto]
    lista = [elemento.replace('\n', '').replace('\t', '') for elemento in lista]
    return lista

# Cargar archivo
def texto_lista(file):
    if file.endswith('.docx'):
        return docx_list(file)
    elif file.endswith('.pdf'):
        return pdf_list(file)
    elif file.endswith('.pptx'):
        return pptx_list(file)
    elif file.endswith('.xlsx'):
        return xlsx_list(file)
    elif file.endswith('.txt'):
        return txt_list(file)
    else:
        return ['Formato no soportado']

st.set_page_config('preguntaDOC')
st.header("Pregunta a tu PDF")
OPENAI_API_KEY = st.text_input('OpenAI API Key', type='password')
pdf_obj = st.file_uploader("Carga tu documento", type="pdf", on_change=st.cache_resource.clear)

@st.cache_resource 
def create_embeddings(file):

    text = texto_lista(file)

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100,
        length_function=len
        )        
    chunks = text_splitter.split_text(text)

    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2")
    knowledge_base = FAISS.from_texts(chunks, embeddings)

    return knowledge_base

if pdf_obj:
    knowledge_base = create_embeddings(pdf_obj)
    user_question = st.text_input("Haz una pregunta sobre tu PDF:")

    if user_question:
        os.environ["OPENAI_API_KEY"] = OPENAI_API_KEY
        docs = knowledge_base.similarity_search(user_question, 3)
        llm = ChatOpenAI(model_name='gpt-3.5-turbo')
        chain = load_qa_chain(llm, chain_type="stuff")
        respuesta = chain.run(input_documents=docs, question=user_question)
        st.write(respuesta)    